# -*- coding: utf-8 -*-
"""Convert the ai_first_course sensor deck markdown files to PPTX with speaker notes."""
import re, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO_ROOT = Path(__file__).resolve().parent.parent
SRC = REPO_ROOT / "sensor_decks"
OUT = REPO_ROOT / "pptx"
OUT.mkdir(exist_ok=True)
EXPECTED = 9

def save_with_retry(prs, path, attempts=4, delay=3):
    import time
    for i in range(attempts):
        try:
            prs.save(path)
            return True
        except PermissionError:
            time.sleep(delay)
    return False

ACCENT = RGBColor(0x1B, 0x5E, 0x8A)
DARK = RGBColor(0x22, 0x22, 0x22)

SLIDE_RE = re.compile(
    r"^## (Slide [\d]+\.\d+: .+?)\n+\*\*Slide content:\*\*\n(.*?)\n+\*\*Narration:\*\*\n(.*?)(?=\n## |\Z)",
    re.S | re.M,
)

def parse_deck(path):
    text = path.read_text(encoding="utf-8")
    m = re.match(r"# (.+)", text)
    deck_title = m.group(1).strip()
    slides = []
    for sm in SLIDE_RE.finditer(text):
        heading = sm.group(1).strip()
        bullets = [b.strip()[2:].strip() for b in sm.group(2).strip().splitlines() if b.strip().startswith("- ")]
        narration = sm.group(3).strip()
        slides.append((heading, bullets, narration))
    return deck_title, slides

def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = ACCENT
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.0))
    tf2 = sub.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(18); r2.font.color.rgb = DARK
    return slide

def add_content_slide(prs, heading, bullets, narration):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # slide number prefix like "Slide 09.3: Technical Card" -> keep full heading as title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = heading
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = ACCENT
    # rule line
    ln = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(0.1))
    # body bullets
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
    btf = body.text_frame; btf.word_wrap = True
    for i, b in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.space_after = Pt(10)
        run = para.add_run(); run.text = "\u2022  " + b
        run.font.size = Pt(19); run.font.color.rgb = DARK
    # speaker notes
    notes = slide.notes_slide.notes_text_frame
    notes.text = narration
    return slide

def build(files):
    combined = new_prs()
    add_title_slide(
        combined,
        "AIoT and Climate Change: Sensor Modules",
        "IEEE BLP / RadioStudio course, AI-first edition. 14 sensors, 5 Sensor Interface Modules. Combined deck with speaker notes.",
    )
    total = 0
    locked = []
    for f in files:
        deck_title, slides = parse_deck(f)
        if len(slides) != EXPECTED:
            print(f"WARNING: {f.name} parsed {len(slides)} slides, expected {EXPECTED}")
        # individual deck
        prs = new_prs()
        add_title_slide(prs, deck_title, "AIoT and Climate Change (AI-first edition). Speaker notes carry the full narration.")
        for heading, bullets, narration in slides:
            add_content_slide(prs, heading, bullets, narration)
        out = OUT / (f.stem + ".pptx")
        if save_with_retry(prs, out):
            print(f"wrote {out.name}: {len(slides)} content slides")
        else:
            locked.append(out.name)
            print(f"LOCKED, not written: {out.name} (close it in PowerPoint and rerun)")
        # combined
        add_title_slide(combined, deck_title, "")
        for heading, bullets, narration in slides:
            add_content_slide(combined, heading, bullets, narration)
        total += len(slides)
    if save_with_retry(combined, OUT / "AIoT_Sensor_Decks_ALL.pptx"):
        print(f"wrote AIoT_Sensor_Decks_ALL.pptx: {total} content slides + 15 title slides")
    else:
        locked.append("AIoT_Sensor_Decks_ALL.pptx")
        print("LOCKED, not written: AIoT_Sensor_Decks_ALL.pptx")
    if locked:
        print("STALE FILES REMAIN:", ", ".join(locked))

if __name__ == "__main__":
    files = sorted(SRC.glob("*.md"))
    build(files)
